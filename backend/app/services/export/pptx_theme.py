"""Shared visual primitives for `PptxExporter`: brand colors, a safe system
font, deterministic vector icons (no external image/font APIs, no emoji glyph
dependence), and small layout helpers (text capacity estimation, card/pill
chrome) reused across every slide builder in `pptx_exporter.py`.

Colors intentionally mirror `_BRAND_COLOR` in `pdf_exporter.py`/
`docx_exporter.py` (`#6D28D9`) plus a small semantic accent set (decisions =
success green, risks = warning amber, questions/timeline = neutral slate) so
slides read at a glance instead of as a wall of uniform bullets.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from app.services.export.content import (
    SECTION_ACTIONS,
    SECTION_DECISIONS,
    SECTION_DISCUSSION,
    SECTION_NEXT_STEPS,
    SECTION_QUESTIONS,
    SECTION_RISKS,
    SECTION_SUMMARY,
    SECTION_TIMELINE,
    SECTION_TRANSCRIPT,
)

# ---------------------------------------------------------------------------
# Palette. `#6D28D9` (violet-700) is the existing PDF/DOCX brand accent; the
# rest extends it with a small, print-safe semantic set rather than inventing
# an unrelated new palette.
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x6D, 0x28, 0xD9)  # violet-700 - brand, summary/discussion/actions/next steps
PRIMARY_DARK = RGBColor(0x3B, 0x0F, 0x6E)  # deep violet - title/appendix bands
PRIMARY_LIGHT = RGBColor(0xEF, 0xE9, 0xFE)  # violet-100 - tint backgrounds

SUCCESS = RGBColor(0x15, 0x80, 0x3D)  # green-700 - decisions
SUCCESS_LIGHT = RGBColor(0xE3, 0xF6, 0xE9)

WARNING = RGBColor(0xB4, 0x53, 0x09)  # amber-700 - risks/blockers
WARNING_LIGHT = RGBColor(0xFC, 0xEF, 0xD9)

NEUTRAL = RGBColor(0x47, 0x55, 0x69)  # slate-600 - questions/timeline/transcript
NEUTRAL_LIGHT = RGBColor(0xEE, 0xF1, 0xF5)

INK = RGBColor(0x11, 0x18, 0x27)  # slate-900 - primary text
SUBTLE = RGBColor(0x5B, 0x67, 0x7A)  # slate-500 - secondary text
MUTED = RGBColor(0x94, 0x9F, 0xB2)  # slate-400 - footer/meta text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)  # slate-200 - hairlines
CANVAS = RGBColor(0xFC, 0xFC, 0xFD)  # near-white slide background

FONT = "Calibri"  # Office-default sans; Geist (web-only) isn't guaranteed on
# the rendering host, and Calibri has metric-compatible fallbacks (Carlito) on
# non-Windows renderers, so it stays legible everywhere without embedding fonts.

_KIND_COLORS: dict[str, tuple[RGBColor, RGBColor]] = {
    SECTION_SUMMARY: (PRIMARY, PRIMARY_LIGHT),
    SECTION_DISCUSSION: (PRIMARY, PRIMARY_LIGHT),
    SECTION_DECISIONS: (SUCCESS, SUCCESS_LIGHT),
    SECTION_ACTIONS: (PRIMARY, PRIMARY_LIGHT),
    SECTION_RISKS: (WARNING, WARNING_LIGHT),
    SECTION_QUESTIONS: (NEUTRAL, NEUTRAL_LIGHT),
    SECTION_NEXT_STEPS: (PRIMARY, PRIMARY_LIGHT),
    SECTION_TIMELINE: (NEUTRAL, NEUTRAL_LIGHT),
    SECTION_TRANSCRIPT: (NEUTRAL, NEUTRAL_LIGHT),
}


def kind_color(kind: str) -> RGBColor:
    return _KIND_COLORS.get(kind, (PRIMARY, PRIMARY_LIGHT))[0]


def kind_tint(kind: str) -> RGBColor:
    return _KIND_COLORS.get(kind, (PRIMARY, PRIMARY_LIGHT))[1]


# ---------------------------------------------------------------------------
# Text capacity estimation. There is no renderer available at generation time
# to measure real glyph metrics, so overflow is prevented by a conservative
# analytic estimate (average proportional-sans char width ~= 0.50em) rather
# than by hoping PowerPoint's shrink-to-fit produces the same result in every
# viewer. Underestimating capacity trades a little whitespace for a hard
# guarantee against clipped text.
# ---------------------------------------------------------------------------
_AVG_CHAR_WIDTH_EM = 0.50
_SAFETY = 0.86


def estimate_capacity(width_in: float, height_in: float, font_pt: float, line_spacing: float = 1.28) -> int:
    chars_per_line = max(1, int((width_in * 72) / (_AVG_CHAR_WIDTH_EM * font_pt)))
    line_height_pt = font_pt * line_spacing
    lines = max(1, int((height_in * 72) / line_height_pt))
    return max(20, int(chars_per_line * lines * _SAFETY))


def truncate_to_fit(text: str, max_chars: int) -> str:
    text = " ".join(text.split())
    if len(text) <= max_chars:
        return text
    cut = text[: max(1, max_chars - 1)]
    last_space = cut.rfind(" ")
    if last_space > max_chars * 0.6:
        cut = cut[:last_space]
    return cut.rstrip(" ,;:-") + "…"


def chunk(items: list, size: int) -> list[list]:
    return [items[i : i + size] for i in range(0, len(items), size)] or [[]]


# ---------------------------------------------------------------------------
# Flowing multi-paragraph text (executive summary, full transcript). A pure
# char-count budget undercounts vertical space when text contains paragraph
# breaks, since a blank line consumes a full line height but zero chars -
# feeding it straight into an overflow-prone slide is exactly the bug this
# guards against. Simulating the actual word-wrap into rendered lines (one
# entry per line, "" for a blank line) keeps the estimate tied to what
# `add_text` will actually lay out, since it also paragraph-splits on "\n".
# ---------------------------------------------------------------------------


def _chars_per_line(width_in: float, font_pt: float) -> int:
    return max(10, int((width_in * 72) / (_AVG_CHAR_WIDTH_EM * font_pt)))


def estimate_line_capacity(height_in: float, font_pt: float, line_spacing: float = 1.3) -> int:
    line_height_pt = font_pt * line_spacing
    return max(3, int((height_in * 72) / line_height_pt * 0.88))


def wrap_lines(text: str, width_in: float, font_pt: float) -> list[str]:
    cpl = _chars_per_line(width_in, font_pt)
    lines: list[str] = []
    for para in text.split("\n"):
        if not para.strip():
            lines.append("")
            continue
        current = ""
        for word in para.split():
            candidate = f"{current} {word}".strip()
            if len(candidate) <= cpl:
                current = candidate
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
    return lines


def _paragraph_line_count(para: str, width_in: float, font_pt: float) -> int:
    return len(wrap_lines(para, width_in, font_pt))


def _split_paragraph_by_lines(para: str, width_in: float, font_pt: float, line_budget: int) -> tuple[str, str]:
    """Returns `(chunk, remainder)`: the longest word-boundary prefix of
    `para` that wraps to at most `line_budget` lines at this width/font.
    """
    words = para.split()
    cpl = _chars_per_line(width_in, font_pt)
    lines_used, current_len = 0, 0
    for index, word in enumerate(words):
        add_len = len(word) if current_len == 0 else len(word) + 1
        if current_len + add_len > cpl:
            lines_used += 1
            current_len = len(word)
            if lines_used >= line_budget:
                return " ".join(words[:index]), " ".join(words[index:])
        else:
            current_len += add_len
    return para, ""


def truncate_to_lines(text: str, width_in: float, font_pt: float, max_lines: int) -> str:
    """Truncates `text` to at most `max_lines` wrapped lines at this
    width/font (word-boundary safe, "…" appended if cut). Uses the same wrap
    simulation a caller uses to size the row/box holding the text, instead of
    the separately-tuned char budget in `truncate_to_fit` - two independent
    estimates can disagree, and a box correctly sized for N lines getting
    text truncated for a *smaller* apparent budget is a needless content loss.
    """
    if len(wrap_lines(text, width_in, font_pt)) <= max_lines:
        return text
    chunk, _ = _split_paragraph_by_lines(text, width_in, font_pt, max_lines)
    return chunk.rstrip(" ,;:-") + "…"


def paginate_flowing_text(
    text: str, width_in: float, height_in: float, font_pt: float, line_spacing: float = 1.3
) -> list[str]:
    """Splits `text` (paragraph breaks as blank lines, i.e. "\\n\\n") into
    page-sized chunks that are each guaranteed to fit `width_in` x
    `height_in` at `font_pt` - by simulating the same word-wrap `add_text`'s
    real rendering will do, so pagination is driven by actual line counts
    instead of a raw character budget (which undercounts the vertical space
    paragraph breaks consume). Returned chunks are real paragraph text
    (paragraphs still joined by "\\n\\n"), not manually pre-wrapped lines -
    the caller's text frame still does its own (matching) word-wrap.
    """
    text = text.strip()
    if not text:
        return [""]
    max_lines = estimate_line_capacity(height_in, font_pt, line_spacing)
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    pages: list[str] = []
    current: list[str] = []
    current_lines = 0

    for para in paragraphs:
        remaining = para
        while remaining:
            # A paragraph joining onto a non-empty page costs one extra
            # rendered line for the blank-line separator ("\n\n" -> an empty
            # paragraph in the text frame) - omitting this was the original
            # bug: five 4-line paragraphs plus their four separators is 24
            # rendered lines, not 20, and the 21st-24th spilled past the slide.
            separator_cost = 1 if current else 0
            budget = max_lines - current_lines - separator_cost if current else max_lines
            if budget <= 0:
                pages.append("\n\n".join(current))
                current, current_lines = [], 0
                separator_cost = 0
                budget = max_lines
            para_lines = _paragraph_line_count(remaining, width_in, font_pt)
            if para_lines <= budget:
                current.append(remaining)
                current_lines += para_lines + separator_cost
                remaining = ""
            else:
                chunk, remaining = _split_paragraph_by_lines(remaining, width_in, font_pt, budget)
                if not chunk:
                    # budget too small for even one word at the start of a
                    # fresh page - force one word through rather than loop.
                    words = remaining.split(" ", 1)
                    chunk, remaining = words[0], (words[1] if len(words) > 1 else "")
                current.append(chunk)
                current_lines += _paragraph_line_count(chunk, width_in, font_pt) + separator_cost
                if remaining:
                    pages.append("\n\n".join(current))
                    current, current_lines = [], 0

    if current:
        pages.append("\n\n".join(current))
    return pages or [""]


# ---------------------------------------------------------------------------
# Low-level shape/text helpers
# ---------------------------------------------------------------------------


def _no_shadow(shape) -> None:
    shape.shadow.inherit = False


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill: RGBColor | None,
    *,
    rounded=False,
    line: RGBColor | None = None,
    line_width_pt: float = 0.75,
    corner_radius: float = 0.06,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        try:
            shape.adjustments[0] = corner_radius
        except IndexError:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width_pt)
    _no_shadow(shape)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: float,
    color: RGBColor = INK,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = FONT,
    line_spacing: float | None = None,
    wrap: bool = True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paragraphs = text.split("\n") or [text]
    for index, para_text in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        if line_spacing:
            # A fixed-point value (not a bare float multiple) so the actual
            # rendered line height exactly matches what `estimate_capacity`/
            # `estimate_line_capacity` assumed - renderers interpret a float
            # multiple against their own (taller, unspecified) baseline line
            # height, which under-filled our capacity math enough to overflow
            # the footer on multi-page text.
            paragraph.line_spacing = Pt(size * line_spacing)
        run = paragraph.add_run()
        run.text = para_text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_pill(
    slide, left, top, width, height, text: str, *, fill: RGBColor, text_color: RGBColor, size: float = 10, bold: bool = True
):
    shape = add_rect(slide, left, top, width, height, fill, rounded=True, corner_radius=0.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = text_color
    return shape


def add_circle(slide, left, top, diameter, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    _no_shadow(shape)
    return shape


def add_dot(slide, left, top, diameter, color: RGBColor):
    """Small solid-circle bullet marker for list rows (as opposed to the
    larger `draw_icon_chip`, which is reserved for slide/panel headers).
    """
    return add_circle(slide, left, top, diameter, color)


def _freeform_check(slide, left, top, size, color: RGBColor, weight_pt: float = 2.25):
    scale = size / 100
    builder = slide.shapes.build_freeform(18, 55, scale=scale)
    builder.add_line_segments([(40, 78), (86, 22)], close=False)
    shape = builder.convert_to_shape(left, top)
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(weight_pt)
    _no_shadow(shape)
    return shape


def draw_icon_chip(slide, kind: str, left, top, diameter=Emu(4200000)) -> None:
    """Draws a small circular icon chip for a section `kind`. Every glyph is
    either a native PowerPoint autoshape or a short freeform mark (or, for
    "?", a plain ASCII character) - no emoji or symbol-font glyphs, so
    rendering is identical across PowerPoint, LibreOffice, Keynote and
    Google Slides.
    """
    color = kind_color(kind)
    chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    _no_shadow(chip)

    diameter_in = Emu(diameter).inches
    inner = Emu(int(diameter * 0.5))
    inner_left = Emu(int(left + (diameter - inner) / 2))
    inner_top = Emu(int(top + (diameter - inner) / 2))

    if kind == SECTION_SUMMARY:
        star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, inner_left, inner_top, inner, inner)
        star.fill.solid()
        star.fill.fore_color.rgb = WHITE
        star.line.fill.background()
        _no_shadow(star)
    elif kind in (SECTION_DISCUSSION, SECTION_TIMELINE, SECTION_TRANSCRIPT):
        bubble_w = Emu(int(diameter * 0.56))
        bubble_h = Emu(int(diameter * 0.44))
        bubble_left = Emu(int(left + (diameter - bubble_w) / 2))
        bubble_top = Emu(int(top + (diameter - bubble_h) / 2 - diameter * 0.04))
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, bubble_left, bubble_top, bubble_w, bubble_h
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = WHITE
        bubble.line.fill.background()
        _no_shadow(bubble)
    elif kind == SECTION_DECISIONS:
        _freeform_check(slide, inner_left, inner_top, inner, WHITE)
    elif kind == SECTION_ACTIONS:
        box_size = Emu(int(diameter * 0.42))
        box_left = Emu(int(left + (diameter - box_size) / 2))
        box_top = Emu(int(top + (diameter - box_size) / 2))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_size, box_size)
        try:
            box.adjustments[0] = 0.2
        except IndexError:
            pass
        box.fill.background()
        box.line.color.rgb = WHITE
        box.line.width = Pt(1.75)
        _no_shadow(box)
    elif kind == SECTION_RISKS:
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, inner_left, inner_top, inner, inner)
        tri.fill.solid()
        tri.fill.fore_color.rgb = WHITE
        tri.line.fill.background()
        _no_shadow(tri)
        add_text(
            slide, left, Emu(int(top + diameter * 0.16)), diameter, Emu(int(diameter * 0.68)),
            "!", size=diameter_in * 26, color=color, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False,
        )
    elif kind == SECTION_NEXT_STEPS:
        arrow_h = Emu(int(inner * 0.62))
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inner_left, inner_top, inner, arrow_h)
        arrow.top = Emu(int(top + (diameter - arrow_h) / 2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = WHITE
        arrow.line.fill.background()
        _no_shadow(arrow)
    elif kind == SECTION_QUESTIONS:
        add_text(
            slide, left, top, diameter, diameter, "?",
            size=diameter_in * 34, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False,
        )
