"""Presentation-oriented PPTX renderer for Meeting Notes.

Unlike `pdf_exporter.py`/`docx_exporter.py` (one flowing document), a
presentation needs *slides*: content grouped into scannable, visually
distinct units instead of one bullet-per-line document poured into slide
placeholders. This module owns that translation from `ExportDocument`
sections into a bounded set of purpose-built slide layouts (title, executive
summary + at-a-glance stats, topic cards, decisions/action items,
risk/question/next-step outlook, discussion timeline, transcript appendix) -
see `pptx_theme.py` for the shared colors, fonts, icons and layout helpers
every builder below uses.

Slide count adapts to content: empty sections never produce a slide, small
sections combine onto one slide, and any section too big for one slide
paginates onto "(cont.)" slides. Nothing here changes what data goes into an
export (see `content.py`) or how PDF/DOCX render it - only how PPTX lays the
same `ExportDocument` out.
"""

from __future__ import annotations

import io
import math

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.export import pptx_theme as theme
from app.services.export.content import (
    SECTION_ACTIONS,
    SECTION_DECISIONS,
    SECTION_DISCUSSION,
    SECTION_NEXT_STEPS,
    SECTION_QUESTIONS,
    SECTION_RISKS,
    SECTION_SUMMARY,
    SECTION_TIMELINE,
    SECTION_TRANSCRIPT,
    ExportDocument,
    ExportSection,
)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.7
CONTENT_WIDTH_IN = SLIDE_W_IN - 2 * MARGIN_IN
HEADER_TOP_IN = 0.5
CHIP_SIZE_IN = 0.5
HEADER_DIVIDER_Y_IN = 1.32
CONTENT_TOP_IN = 1.55
CONTENT_BOTTOM_IN = 6.85
CONTENT_HEIGHT_IN = CONTENT_BOTTOM_IN - CONTENT_TOP_IN
FOOTER_DIVIDER_Y_IN = 6.98
_BLANK_LAYOUT = 6
_HAIRLINE = Emu(9525)  # ~0.75pt

_TRANSCRIPT_MIN_CHARS = 300
_TOPIC_MAX_PER_SLIDE = 6
_DECISION_ROWS_PER_SLIDE = 8
_ACTION_ROWS_COMBINED_MAX = 6
_ACTION_ROWS_PER_SLIDE = 9
_OUTLOOK_ROWS_COMBINED_MAX = 6
_OUTLOOK_ROWS_PER_SLIDE = 8
_TIMELINE_ROWS_PER_SLIDE = 7

_STAT_ORDER = [
    SECTION_DISCUSSION,
    SECTION_DECISIONS,
    SECTION_ACTIONS,
    SECTION_RISKS,
    SECTION_QUESTIONS,
    SECTION_NEXT_STEPS,
]
_STAT_LABELS = {
    SECTION_DISCUSSION: "Discussion Topics",
    SECTION_DECISIONS: "Decisions",
    SECTION_ACTIONS: "Action Items",
    SECTION_RISKS: "Risks / Blockers",
    SECTION_QUESTIONS: "Open Questions",
    SECTION_NEXT_STEPS: "Next Steps",
}
_STATUS_COLORS = {
    "done": theme.SUCCESS,
    "completed": theme.SUCCESS,
    "complete": theme.SUCCESS,
    "in progress": theme.PRIMARY,
    "in-progress": theme.PRIMARY,
    "blocked": theme.WARNING,
    "at risk": theme.WARNING,
    "not started": theme.NEUTRAL,
}


def _in(value: float) -> Emu:
    return Inches(value)


class PptxExporter:
    content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    file_extension = "pptx"

    def render(self, document: ExportDocument) -> bytes:
        prs = Presentation()
        prs.slide_width = _in(SLIDE_W_IN)
        prs.slide_height = _in(SLIDE_H_IN)
        by_kind = {section.kind: section for section in document.sections}

        self._title_slide(prs, document)

        summary = by_kind.get(SECTION_SUMMARY)
        if summary:
            self._summary_slides(prs, document, summary, by_kind)

        discussion = by_kind.get(SECTION_DISCUSSION)
        if discussion:
            self._topic_slides(prs, document, discussion)

        decisions = by_kind.get(SECTION_DECISIONS)
        actions = by_kind.get(SECTION_ACTIONS)
        if decisions or actions:
            self._decisions_actions_slides(prs, document, decisions, actions)

        risks = by_kind.get(SECTION_RISKS)
        questions = by_kind.get(SECTION_QUESTIONS)
        next_steps = by_kind.get(SECTION_NEXT_STEPS)
        if risks or questions or next_steps:
            self._outlook_slides(prs, document, risks, questions, next_steps)

        timeline = by_kind.get(SECTION_TIMELINE)
        if timeline:
            self._timeline_slides(prs, document, timeline)

        transcript = by_kind.get(SECTION_TRANSCRIPT)
        if transcript and len(transcript.lines[0].strip()) >= _TRANSCRIPT_MIN_CHARS:
            self._appendix_divider(prs)
            self._transcript_slides(prs, document, transcript)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    # -- chrome ------------------------------------------------------------

    def _blank_slide(self, prs: Presentation):
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
        theme.add_rect(slide, _in(0), _in(0), _in(SLIDE_W_IN), _in(SLIDE_H_IN), theme.CANVAS)
        return slide

    def _footer(self, slide, document: ExportDocument, page_label: str | None = None) -> None:
        theme.add_rect(slide, _in(MARGIN_IN), _in(FOOTER_DIVIDER_Y_IN), _in(CONTENT_WIDTH_IN), _HAIRLINE, theme.BORDER)
        theme.add_text(
            slide, _in(MARGIN_IN), _in(FOOTER_DIVIDER_Y_IN + 0.08), _in(3.5), _in(0.26),
            document.brand.upper(), size=8.5, bold=True, color=theme.PRIMARY,
        )
        right_text = document.meeting_title if not page_label else f"{document.meeting_title} · {page_label}"
        right_text = theme.truncate_to_fit(right_text, 95)
        theme.add_text(
            slide, _in(MARGIN_IN + 3.5), _in(FOOTER_DIVIDER_Y_IN + 0.08), _in(CONTENT_WIDTH_IN - 3.5), _in(0.26),
            right_text, size=8.5, color=theme.MUTED, align=PP_ALIGN.RIGHT,
        )

    def _content_slide(
        self, prs: Presentation, document: ExportDocument, kind: str, heading: str,
        subtitle: str | None = None, page_label: str | None = None,
    ):
        slide = self._blank_slide(prs)
        theme.draw_icon_chip(slide, kind, _in(MARGIN_IN), _in(HEADER_TOP_IN), _in(CHIP_SIZE_IN))
        text_left = MARGIN_IN + CHIP_SIZE_IN + 0.22
        text_width = CONTENT_WIDTH_IN - CHIP_SIZE_IN - 0.22
        theme.add_text(
            slide, _in(text_left), _in(HEADER_TOP_IN - 0.04), _in(text_width), _in(0.42),
            heading, size=22, bold=True, color=theme.INK, anchor=MSO_ANCHOR.MIDDLE,
        )
        if subtitle:
            theme.add_text(
                slide, _in(text_left), _in(HEADER_TOP_IN + 0.37), _in(text_width), _in(0.24),
                subtitle, size=11, color=theme.SUBTLE,
            )
        theme.add_rect(slide, _in(MARGIN_IN), _in(HEADER_DIVIDER_Y_IN), _in(CONTENT_WIDTH_IN), _HAIRLINE, theme.BORDER)
        self._footer(slide, document, page_label)
        return slide

    def _plain_content_slide(self, prs: Presentation, document: ExportDocument, heading: str, page_label: str | None = None):
        """Header without a single section icon, for slides that combine more
        than one section kind (each sub-panel draws its own small icon).
        """
        slide = self._blank_slide(prs)
        theme.add_text(
            slide, _in(MARGIN_IN), _in(HEADER_TOP_IN - 0.04), _in(CONTENT_WIDTH_IN), _in(0.42),
            heading, size=22, bold=True, color=theme.INK, anchor=MSO_ANCHOR.MIDDLE,
        )
        theme.add_rect(slide, _in(MARGIN_IN), _in(HEADER_DIVIDER_Y_IN), _in(CONTENT_WIDTH_IN), _HAIRLINE, theme.BORDER)
        self._footer(slide, document, page_label)
        return slide

    # -- layout ---------------------------------------------------------------

    @staticmethod
    def _centered_top(natural_height: float, top_in: float, height_in: float, max_offset: float = 1.6) -> float:
        """Nudges a content block down from `top_in` when it's much shorter
        than its slide's content area, so a handful of decisions/actions/
        risks doesn't read as a mostly-empty slide - without fully centering,
        which would break the block's visual association with its heading.
        """
        slack = height_in - natural_height
        if slack <= 0.6:
            return top_in
        return top_in + min(slack * 0.42, max_offset)

    # -- slide 1: title ------------------------------------------------------

    def _title_slide(self, prs: Presentation, document: ExportDocument) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
        theme.add_rect(slide, _in(0), _in(0), _in(SLIDE_W_IN), _in(SLIDE_H_IN), theme.WHITE)

        band_w = 4.6
        theme.add_rect(slide, _in(0), _in(0), _in(band_w), _in(SLIDE_H_IN), theme.PRIMARY_DARK)
        theme.add_circle(slide, _in(-1.1), _in(SLIDE_H_IN - 2.6), _in(3.4), theme.PRIMARY)

        theme.add_text(slide, _in(0.55), _in(0.55), _in(band_w - 1), _in(0.4), document.brand.upper(), size=15, bold=True, color=theme.WHITE)
        theme.add_text(slide, _in(0.55), _in(0.95), _in(band_w - 1), _in(0.3), "Meeting Notes", size=11, color=theme.PRIMARY_LIGHT)

        right_left = band_w + 0.5
        right_width = SLIDE_W_IN - right_left - MARGIN_IN
        title = theme.truncate_to_fit(document.meeting_title, 130)
        theme.add_text(
            slide, _in(right_left), _in(2.1), _in(right_width), _in(1.9), title,
            size=34, bold=True, color=theme.INK, line_spacing=1.12, anchor=MSO_ANCHOR.BOTTOM,
        )

        chips = []
        if document.date_time_ist:
            chips.append(("Date & Time", document.date_time_ist))
        if document.duration_label:
            chips.append(("Duration", document.duration_label))
        if document.participants_label:
            chips.append(("Participants", document.participants_label))

        if chips:
            gap = 0.25
            chip_w = (right_width - gap * (len(chips) - 1)) / len(chips)
            chip_h = 1.0
            chip_top = 4.35
            for index, (label, value) in enumerate(chips):
                x = right_left + index * (chip_w + gap)
                theme.add_rect(slide, _in(x), _in(chip_top), _in(chip_w), _in(chip_h), theme.PRIMARY_LIGHT, rounded=True)
                theme.add_rect(slide, _in(x), _in(chip_top), _in(0.05), _in(chip_h), theme.PRIMARY)
                theme.add_text(slide, _in(x + 0.22), _in(chip_top + 0.16), _in(chip_w - 0.4), _in(0.24), label.upper(), size=9, bold=True, color=theme.SUBTLE)
                value_capacity = theme.estimate_capacity(chip_w - 0.4, 0.5, 13.5, line_spacing=1.1)
                value_text = theme.truncate_to_fit(value, value_capacity)
                theme.add_text(slide, _in(x + 0.22), _in(chip_top + 0.44), _in(chip_w - 0.4), _in(0.5), value_text, size=13.5, bold=True, color=theme.INK, line_spacing=1.1)

        disclaimer_top = chip_top + 1.0 + 0.25 if chips else 5.6
        theme.add_text(
            slide, _in(right_left), _in(disclaimer_top), _in(right_width), _in(0.4),
            document.disclaimer, size=8.5, italic=True, color=theme.SUBTLE, line_spacing=1.2,
        )

    # -- slide 2: executive summary + at-a-glance ----------------------------

    def _summary_slides(self, prs: Presentation, document: ExportDocument, summary: ExportSection, by_kind: dict[str, ExportSection]) -> None:
        text = "\n\n".join(part.strip() for part in summary.lines[0].split("\n") if part.strip())
        main_font = 14.5
        main_width = CONTENT_WIDTH_IN * 0.60
        gap = 0.5
        side_width = CONTENT_WIDTH_IN - main_width - gap

        pages = theme.paginate_flowing_text(text, main_width - 0.6, CONTENT_HEIGHT_IN - 0.6, main_font, line_spacing=1.32)

        slide = self._content_slide(prs, document, SECTION_SUMMARY, "Executive Summary")
        theme.add_rect(slide, _in(MARGIN_IN), _in(CONTENT_TOP_IN), _in(main_width), _in(CONTENT_HEIGHT_IN), theme.WHITE, rounded=True, line=theme.BORDER, line_width_pt=1)
        theme.add_text(
            slide, _in(MARGIN_IN + 0.3), _in(CONTENT_TOP_IN + 0.3), _in(main_width - 0.6), _in(CONTENT_HEIGHT_IN - 0.6),
            pages[0], size=main_font, color=theme.INK, line_spacing=1.32,
        )

        side_left = MARGIN_IN + main_width + gap
        theme.add_text(slide, _in(side_left), _in(CONTENT_TOP_IN), _in(side_width), _in(0.3), "AT A GLANCE", size=10.5, bold=True, color=theme.SUBTLE)
        stats = self._collect_stats(by_kind)
        self._render_stat_grid(slide, stats, side_left, CONTENT_TOP_IN + 0.4, side_width, CONTENT_HEIGHT_IN - 0.4)

        if len(pages) > 1:
            remaining_text = "\n\n".join(pages[1:])
            cont_pages = theme.paginate_flowing_text(
                remaining_text, CONTENT_WIDTH_IN - 0.6, CONTENT_HEIGHT_IN - 0.6, main_font, line_spacing=1.35
            )
            for extra_page in cont_pages:
                cslide = self._content_slide(prs, document, SECTION_SUMMARY, "Executive Summary (cont.)")
                theme.add_text(
                    cslide, _in(MARGIN_IN + 0.3), _in(CONTENT_TOP_IN + 0.3), _in(CONTENT_WIDTH_IN - 0.6), _in(CONTENT_HEIGHT_IN - 0.6),
                    extra_page, size=main_font, color=theme.INK, line_spacing=1.35,
                )

    def _collect_stats(self, by_kind: dict[str, ExportSection]) -> list[tuple[str, int, str]]:
        stats = []
        for kind in _STAT_ORDER:
            section = by_kind.get(kind)
            if section:
                count = len(section.items) if section.items is not None else len(section.lines)
                stats.append((kind, count, _STAT_LABELS[kind]))
        return stats

    def _render_stat_grid(self, slide, stats: list[tuple[str, int, str]], left_in: float, top_in: float, width_in: float, height_in: float) -> None:
        if not stats:
            theme.add_text(slide, _in(left_in), _in(top_in), _in(width_in), _in(0.4), "No additional sections recorded.", size=11, italic=True, color=theme.SUBTLE)
            return
        cols = 2 if len(stats) > 1 else 1
        rows = math.ceil(len(stats) / cols)
        gap = 0.14
        tile_w = (width_in - gap * (cols - 1)) / cols
        tile_h = min(1.1, (height_in - gap * (rows - 1)) / rows)
        for index, (kind, count, label) in enumerate(stats):
            row, col = divmod(index, cols)
            x = left_in + col * (tile_w + gap)
            y = top_in + row * (tile_h + gap)
            theme.add_rect(slide, _in(x), _in(y), _in(tile_w), _in(tile_h), theme.kind_tint(kind), rounded=True)
            theme.add_rect(slide, _in(x), _in(y), _in(0.06), _in(tile_h), theme.kind_color(kind))
            theme.add_text(slide, _in(x + 0.2), _in(y + 0.12), _in(tile_w - 0.34), _in(0.44), str(count), size=22, bold=True, color=theme.kind_color(kind))
            theme.add_text(slide, _in(x + 0.2), _in(y + tile_h - 0.36), _in(tile_w - 0.34), _in(0.3), label, size=9.5, color=theme.SUBTLE)

    # -- slide 3: discussion topics -----------------------------------------

    def _topic_slides(self, prs: Presentation, document: ExportDocument, section: ExportSection) -> None:
        topics = section.items or [{"title": line, "description": None} for line in section.lines]
        pages = theme.chunk(topics, _TOPIC_MAX_PER_SLIDE)
        total_pages = len(pages)
        for index, page in enumerate(pages):
            heading = "Discussion Topics" if index == 0 else "Discussion Topics (cont.)"
            subtitle = f"{len(topics)} topic{'s' if len(topics) != 1 else ''} discussed" if index == 0 else None
            page_label = f"Topics {index + 1}/{total_pages}" if total_pages > 1 else None
            slide = self._content_slide(prs, document, SECTION_DISCUSSION, heading, subtitle, page_label)
            self._render_topic_cards(slide, page)

    def _render_topic_cards(self, slide, topics: list[dict]) -> None:
        count = len(topics)
        cols = 1 if count == 1 else (2 if count <= 4 else 3)
        rows = math.ceil(count / cols)
        gap = 0.25
        card_w = (CONTENT_WIDTH_IN - gap * (cols - 1)) / cols
        card_h = min(2.3, (CONTENT_HEIGHT_IN - gap * (rows - 1)) / rows)
        title_font = 14 if cols <= 2 else 12.5
        desc_font = 11 if cols <= 2 else 10.5

        grid_top = self._centered_top(rows * card_h + gap * (rows - 1), CONTENT_TOP_IN, CONTENT_HEIGHT_IN, max_offset=1.2)
        for index, topic in enumerate(topics):
            row, col = divmod(index, cols)
            x = MARGIN_IN + col * (card_w + gap)
            y = grid_top + row * (card_h + gap)
            theme.add_rect(slide, _in(x), _in(y), _in(card_w), _in(card_h), theme.WHITE, rounded=True, line=theme.BORDER, line_width_pt=1)
            chip_d = 0.36
            theme.draw_icon_chip(slide, SECTION_DISCUSSION, _in(x + 0.2), _in(y + 0.2), _in(chip_d))
            title = theme.truncate_to_lines(topic.get("title") or "", card_w - 0.9, title_font, 2)
            theme.add_text(
                slide, _in(x + 0.68), _in(y + 0.2), _in(card_w - 0.9), _in(0.5),
                title, size=title_font, bold=True, color=theme.INK, line_spacing=1.15,
            )
            description = (topic.get("description") or "").strip()
            if description and card_h > 1.1:
                desc_top = y + 0.8
                desc_height = card_h - 1.0
                desc_max_lines = max(1, int(desc_height / (desc_font * 1.28 / 72)))
                desc_text = theme.truncate_to_lines(description, card_w - 0.4, desc_font, desc_max_lines)
                theme.add_text(
                    slide, _in(x + 0.2), _in(desc_top), _in(card_w - 0.4), _in(desc_height),
                    desc_text, size=desc_font, color=theme.SUBTLE, line_spacing=1.28,
                )

    # -- slide 4: decisions + action items -----------------------------------

    def _decisions_actions_slides(
        self, prs: Presentation, document: ExportDocument, decisions: ExportSection | None, actions: ExportSection | None
    ) -> None:
        action_items = actions.items if actions else []
        decision_lines = decisions.lines if decisions else []

        combined_ok = (
            decisions is not None
            and actions is not None
            and len(decision_lines) <= _ACTION_ROWS_COMBINED_MAX
            and len(action_items) <= _ACTION_ROWS_COMBINED_MAX
        )

        if combined_ok:
            slide = self._plain_content_slide(prs, document, "Decisions & Action Items")
            half_w = (CONTENT_WIDTH_IN - 0.5) / 2
            self._render_decision_panel(slide, decision_lines, MARGIN_IN, CONTENT_TOP_IN, half_w, CONTENT_HEIGHT_IN, heading="Decisions")
            self._render_action_table(slide, action_items, MARGIN_IN + half_w + 0.5, CONTENT_TOP_IN, half_w, CONTENT_HEIGHT_IN, heading=actions.heading)
            return

        if decisions:
            pages = theme.chunk(decision_lines, _DECISION_ROWS_PER_SLIDE)
            total_pages = len(pages)
            for index, page in enumerate(pages):
                heading = "Decisions" if index == 0 else "Decisions (cont.)"
                page_label = f"Decisions {index + 1}/{total_pages}" if total_pages > 1 else None
                slide = self._content_slide(prs, document, SECTION_DECISIONS, heading, page_label=page_label)
                self._render_decision_panel(slide, page, MARGIN_IN, CONTENT_TOP_IN, CONTENT_WIDTH_IN, CONTENT_HEIGHT_IN)

        if actions:
            pages = theme.chunk(action_items, _ACTION_ROWS_PER_SLIDE)
            total_pages = len(pages)
            for index, page in enumerate(pages):
                heading = actions.heading if index == 0 else f"{actions.heading} (cont.)"
                page_label = f"Actions {index + 1}/{total_pages}" if total_pages > 1 else None
                slide = self._content_slide(prs, document, SECTION_ACTIONS, heading, page_label=page_label)
                self._render_action_table(slide, page, MARGIN_IN, CONTENT_TOP_IN, CONTENT_WIDTH_IN, CONTENT_HEIGHT_IN)

    _DECISION_MAX_LINES = 2

    def _render_decision_panel(
        self, slide, items: list[str], left_in: float, top_in: float, width_in: float, height_in: float, heading: str | None = None
    ) -> None:
        header_h = 0.55 if heading else 0.0
        font_size = 12
        text_width = width_in - 0.5
        line_h_in = font_size * 1.15 / 72

        row_heights = [
            max(0.42, min(len(theme.wrap_lines(text, text_width, font_size)), self._DECISION_MAX_LINES) * line_h_in + 0.16)
            for text in items
        ]
        avail = height_in - header_h
        rows_total = sum(row_heights)
        if items and rows_total > avail:
            scale = avail / rows_total
            row_heights = [h * scale for h in row_heights]
            rows_total = avail

        natural = header_h + (rows_total if items else 0.4)
        y = self._centered_top(natural, top_in, height_in)

        if heading:
            theme.draw_icon_chip(slide, SECTION_DECISIONS, _in(left_in), _in(y), _in(0.34))
            theme.add_text(slide, _in(left_in + 0.46), _in(y + 0.03), _in(width_in - 0.46), _in(0.3), heading, size=13, bold=True, color=theme.INK)
            y += header_h

        if not items:
            theme.add_text(slide, _in(left_in), _in(y), _in(width_in), _in(0.4), "No decisions recorded.", size=11, italic=True, color=theme.SUBTLE)
            return

        ry = y
        for index, (text, row_h) in enumerate(zip(items, row_heights)):
            if index % 2 == 1:
                theme.add_rect(slide, _in(left_in), _in(ry), _in(width_in), _in(row_h), theme.SUCCESS_LIGHT)
            theme.add_dot(slide, _in(left_in + 0.12), _in(ry + row_h / 2 - 0.045), _in(0.09), theme.SUCCESS)
            line = theme.truncate_to_lines(text, text_width, font_size, self._DECISION_MAX_LINES)
            theme.add_text(
                slide, _in(left_in + 0.32), _in(ry), _in(width_in - 0.44), _in(row_h),
                line, size=font_size, color=theme.INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15,
            )
            ry += row_h

    def _render_action_table(
        self, slide, items: list[dict], left_in: float, top_in: float, width_in: float, height_in: float, heading: str | None = None
    ) -> None:
        header_h = 0.55 if heading else 0.0
        rows = len(items) + 1
        table_h = min(height_in - header_h, 0.45 * rows) if items else 0.4
        natural = header_h + table_h
        y = self._centered_top(natural, top_in, height_in)

        if heading:
            theme.draw_icon_chip(slide, SECTION_ACTIONS, _in(left_in), _in(y), _in(0.34))
            theme.add_text(slide, _in(left_in + 0.46), _in(y + 0.03), _in(width_in - 0.46), _in(0.3), heading, size=13, bold=True, color=theme.INK)
            y += header_h

        if not items:
            theme.add_text(slide, _in(left_in), _in(y), _in(width_in), _in(0.4), "No action items recorded.", size=11, italic=True, color=theme.SUBTLE)
            return
        shape = slide.shapes.add_table(rows, 4, _in(left_in), _in(y), _in(width_in), _in(table_h))
        table = shape.table
        table.first_row = False

        col_fractions = [0.50, 0.18, 0.16, 0.16]
        for col_index, frac in enumerate(col_fractions):
            table.columns[col_index].width = _in(width_in * frac)

        for col_index, label in enumerate(["Action", "Owner", "Due", "Status"]):
            self._style_cell(table.cell(0, col_index), label, fill=theme.PRIMARY, color=theme.WHITE, bold=True, size=10.5)

        row_h_in = table_h / rows
        action_col_width = width_in * col_fractions[0] - 0.3
        action_max_lines = max(1, int((row_h_in - 0.06) / (10.5 * 1.15 / 72)))
        for row_index, item in enumerate(items, start=1):
            row_fill = theme.WHITE if row_index % 2 else theme.NEUTRAL_LIGHT
            action_text = theme.truncate_to_lines(item.get("text") or "", action_col_width, 10.5, action_max_lines)
            self._style_cell(table.cell(row_index, 0), action_text, fill=row_fill, color=theme.INK, size=10.5)
            self._style_cell(table.cell(row_index, 1), item.get("owner") or "—", fill=row_fill, color=theme.INK, size=10.5)
            self._style_cell(table.cell(row_index, 2), item.get("due_date") or "—", fill=row_fill, color=theme.INK, size=10.5)
            status_color = _STATUS_COLORS.get((item.get("status") or "").strip().lower(), theme.NEUTRAL)
            self._style_cell(table.cell(row_index, 3), item.get("status") or "—", fill=row_fill, color=status_color, bold=True, size=10.5)

    @staticmethod
    def _style_cell(cell, text: str, *, fill, color, bold: bool = False, size: float = 10.5) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Emu(45720)
        cell.margin_top = cell.margin_bottom = Emu(13716)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        paragraph = tf.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = theme.FONT
        run.font.color.rgb = color

    # -- slide 5: risks / open questions / next steps ------------------------

    def _outlook_slides(
        self, prs: Presentation, document: ExportDocument,
        risks: ExportSection | None, questions: ExportSection | None, next_steps: ExportSection | None,
    ) -> None:
        present = [section for section in (risks, questions, next_steps) if section]
        max_items = max((len(section.lines) for section in present), default=0)

        if max_items <= _OUTLOOK_ROWS_COMBINED_MAX:
            heading = present[0].heading if len(present) == 1 else "Risks, Questions & Next Steps"
            slide = self._plain_content_slide(prs, document, heading)
            gap = 0.4
            col_w = (CONTENT_WIDTH_IN - gap * (len(present) - 1)) / len(present)
            for index, section in enumerate(present):
                x = MARGIN_IN + index * (col_w + gap)
                self._render_bullet_column(slide, section, x, CONTENT_TOP_IN, col_w, CONTENT_HEIGHT_IN)
            return

        for section in present:
            pages = theme.chunk(section.lines, _OUTLOOK_ROWS_PER_SLIDE)
            total_pages = len(pages)
            for index, page in enumerate(pages):
                heading = section.heading if index == 0 else f"{section.heading} (cont.)"
                page_label = f"{section.heading} {index + 1}/{total_pages}" if total_pages > 1 else None
                slide = self._content_slide(prs, document, section.kind, heading, page_label=page_label)
                self._render_bullet_column(
                    slide, ExportSection(section.heading, page, kind=section.kind),
                    MARGIN_IN, CONTENT_TOP_IN, CONTENT_WIDTH_IN, CONTENT_HEIGHT_IN, compact=False,
                )

    _BULLET_MAX_LINES = 3

    def _render_bullet_column(
        self, slide, section: ExportSection, left_in: float, top_in: float, width_in: float, height_in: float, compact: bool = True
    ) -> None:
        header_h = 0.55
        items = section.lines
        font_size = 12 if compact else 13
        text_width = width_in - 0.32
        line_h_in = font_size * 1.22 / 72

        # Row height follows each item's own wrapped line count (capped) so a
        # short bullet doesn't reserve the same box as a long one, and a long
        # one isn't force-truncated just because its neighbors are short.
        row_heights = [
            max(0.4, min(len(theme.wrap_lines(text, text_width, font_size)), self._BULLET_MAX_LINES) * line_h_in + 0.16)
            for text in items
        ]
        avail = height_in - header_h
        rows_total = sum(row_heights)
        if items and rows_total > avail:
            scale = avail / rows_total
            row_heights = [h * scale for h in row_heights]
            rows_total = avail

        natural = header_h + (rows_total if items else 0.4)
        y = self._centered_top(natural, top_in, height_in, max_offset=1.1)

        theme.draw_icon_chip(slide, section.kind, _in(left_in), _in(y), _in(0.36))
        theme.add_text(
            slide, _in(left_in + 0.48), _in(y + 0.03), _in(width_in - 0.48), _in(0.34),
            section.heading, size=14 if compact else 16, bold=True, color=theme.INK,
        )
        y += header_h
        if not items:
            theme.add_text(slide, _in(left_in), _in(y), _in(width_in), _in(0.4), "None recorded.", size=11, italic=True, color=theme.SUBTLE)
            return

        color = theme.kind_color(section.kind)
        ry = y
        for text, row_h in zip(items, row_heights):
            theme.add_dot(slide, _in(left_in + 0.02), _in(ry + 0.13), _in(0.09), color)
            line = theme.truncate_to_lines(text, text_width, font_size, self._BULLET_MAX_LINES)
            theme.add_text(
                slide, _in(left_in + 0.24), _in(ry), _in(width_in - 0.24), _in(row_h),
                line, size=font_size, color=theme.INK, line_spacing=1.22,
            )
            ry += row_h

    # -- slide 6: detailed discussion timeline -------------------------------

    def _timeline_slides(self, prs: Presentation, document: ExportDocument, section: ExportSection) -> None:
        entries = section.items or [{"timestamp": "", "text": line} for line in section.lines]
        pages = theme.chunk(entries, _TIMELINE_ROWS_PER_SLIDE)
        total_pages = len(pages)
        for index, page in enumerate(pages):
            heading = "Detailed Discussion" if index == 0 else "Detailed Discussion (cont.)"
            page_label = f"Timeline {index + 1}/{total_pages}" if total_pages > 1 else None
            slide = self._content_slide(prs, document, SECTION_TIMELINE, heading, page_label=page_label)
            self._render_timeline(slide, page)

    def _render_timeline(self, slide, entries: list[dict]) -> None:
        if not entries:
            return
        row_h = min(1.0, CONTENT_HEIGHT_IN / len(entries))
        pill_w = 0.85
        rail_x = MARGIN_IN + pill_w + 0.3
        text_left = rail_x + 0.3
        theme.add_rect(slide, _in(rail_x - 0.012), _in(CONTENT_TOP_IN + 0.08), _in(0.024), _in(row_h * len(entries) - 0.16), theme.BORDER)
        for index, entry in enumerate(entries):
            ry = CONTENT_TOP_IN + index * row_h
            theme.add_pill(
                slide, _in(MARGIN_IN), _in(ry + row_h / 2 - 0.15), _in(pill_w), _in(0.3),
                entry.get("timestamp") or "", fill=theme.NEUTRAL_LIGHT, text_color=theme.NEUTRAL, size=9.5,
            )
            theme.add_dot(slide, _in(rail_x - 0.06), _in(ry + row_h / 2 - 0.06), _in(0.12), theme.NEUTRAL)
            text_width = CONTENT_WIDTH_IN - (text_left - MARGIN_IN)
            max_lines = max(1, int((row_h - 0.1) / (11.5 * 1.2 / 72)))
            text = theme.truncate_to_lines(entry.get("text") or "", text_width, 11.5, max_lines)
            theme.add_text(
                slide, _in(text_left), _in(ry), _in(text_width), _in(row_h),
                text, size=11.5, color=theme.INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2,
            )

    # -- appendix: full transcript --------------------------------------------

    def _appendix_divider(self, prs: Presentation) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
        theme.add_rect(slide, _in(0), _in(0), _in(SLIDE_W_IN), _in(SLIDE_H_IN), theme.PRIMARY_DARK)
        theme.add_text(slide, _in(1), _in(3.05), _in(SLIDE_W_IN - 2), _in(0.4), "APPENDIX", size=15, bold=True, color=theme.PRIMARY_LIGHT, align=PP_ALIGN.CENTER)
        theme.add_text(slide, _in(1), _in(3.5), _in(SLIDE_W_IN - 2), _in(0.9), "Full Transcript", size=32, bold=True, color=theme.WHITE, align=PP_ALIGN.CENTER)

    def _transcript_slides(self, prs: Presentation, document: ExportDocument, section: ExportSection) -> None:
        text = "\n\n".join(part.strip() for part in section.lines[0].split("\n") if part.strip()) or section.lines[0].strip()
        font_size = 11.5
        pages = theme.paginate_flowing_text(text, CONTENT_WIDTH_IN - 0.4, CONTENT_HEIGHT_IN - 0.3, font_size, line_spacing=1.32)
        total_pages = len(pages)
        for index, piece in enumerate(pages):
            heading = "Full Transcript" if index == 0 else "Full Transcript (cont.)"
            page_label = f"Transcript {index + 1}/{total_pages}" if total_pages > 1 else None
            slide = self._content_slide(prs, document, SECTION_TRANSCRIPT, heading, page_label=page_label)
            theme.add_text(
                slide, _in(MARGIN_IN + 0.2), _in(CONTENT_TOP_IN + 0.15), _in(CONTENT_WIDTH_IN - 0.4), _in(CONTENT_HEIGHT_IN - 0.3),
                piece, size=font_size, color=theme.INK, line_spacing=1.32,
            )
